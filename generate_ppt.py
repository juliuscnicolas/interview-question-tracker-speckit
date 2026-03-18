"""Generate a 10-slide PowerPoint presentation summarizing the Interview Question Tracker implementation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color palette
PRIMARY = RGBColor(0x1E, 0x40, 0xAF)       # Blue-800
ACCENT = RGBColor(0x3B, 0x82, 0xF6)        # Blue-500
DARK = RGBColor(0x1E, 0x29, 0x3B)          # Slate-800
MEDIUM = RGBColor(0x47, 0x55, 0x69)        # Slate-600
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)     # Slate-100
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)         # Emerald-600
ORANGE = RGBColor(0xEA, 0x58, 0x0C)        # Orange-600
PURPLE = RGBColor(0x7C, 0x3A, 0xED)        # Violet-600
CYAN = RGBColor(0x06, 0x91, 0xB7)          # Cyan-600


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    """Add a colored rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a textbox with text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK, spacing=Pt(8)):
    """Add a bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return txBox


# ─── Slide 1: Title ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, PRIMARY)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), PRIMARY)
add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
            "Interview Question Tracker", 48, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
            "A Vue 3 Single-Page Application for Interview Preparation", 24, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.8),
            "Vue 3  •  Vite  •  Tailwind CSS  •  localStorage  •  GitHub Pages", 18, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.6),
            "Implementation Summary  |  March 2026", 16, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)


# ─── Slide 2: SpecKit Process ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "SpecKit: From Idea to Implementation", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
            "A structured, AI-assisted workflow that turns a one-line idea into production-ready code", 16, MEDIUM)

# Pipeline steps
speckit_steps = [
    ("1", "specify", "spec.md", "Natural language → feature spec\nwith user stories, acceptance\ncriteria & priorities (P1-P3)", RGBColor(0xDC, 0x26, 0x26)),
    ("2", "clarify", "checklist", "Validate spec quality, identify\ngaps, ensure testability &\nno implementation leaks", ORANGE),
    ("3", "plan", "plan.md", "Tech stack decisions, project\nstructure, component design,\nresearch.md + data-model.md", ACCENT),
    ("4", "tasks", "tasks.md", "Dependency-ordered task list\ngrouped by user story with\nparallel execution markers", GREEN),
    ("5", "implement", "src/", "Execute tasks sequentially,\nwrite code, tests, deploy\nconfig — all 34 tasks done", PURPLE),
]

for i, (num, cmd, output, desc, color) in enumerate(speckit_steps):
    x = Inches(0.4 + i * 2.5)
    # Step card
    add_shape(slide, x, Inches(2.2), Inches(2.3), Inches(4.6), LIGHT_BG, 0.03)
    # Number circle
    add_shape(slide, x + Inches(0.8), Inches(2.35), Inches(0.7), Inches(0.7), color, 0.5)
    add_textbox(slide, x + Inches(0.8), Inches(2.4), Inches(0.7), Inches(0.6), num, 24, WHITE, True, PP_ALIGN.CENTER)
    # Command name
    add_textbox(slide, x + Inches(0.15), Inches(3.2), Inches(2.0), Inches(0.5),
                f"/speckit.{cmd}", 14, color, True, PP_ALIGN.CENTER)
    # Output artifact
    add_shape(slide, x + Inches(0.3), Inches(3.7), Inches(1.7), Inches(0.45), color, 0.1)
    add_textbox(slide, x + Inches(0.3), Inches(3.72), Inches(1.7), Inches(0.4),
                output, 13, WHITE, True, PP_ALIGN.CENTER)
    # Description
    txBox = slide.shapes.add_textbox(x + Inches(0.15), Inches(4.4), Inches(2.0), Inches(2.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line_idx, line in enumerate(desc.split("\n")):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

# Arrow connectors between steps
for i in range(4):
    x = Inches(2.8 + i * 2.5)
    add_textbox(slide, x, Inches(3.78), Inches(0.4), Inches(0.4), "→", 20, MEDIUM, True, PP_ALIGN.CENTER)

# Bottom summary
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.7), PRIMARY, 0.05)
add_textbox(slide, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.6),
            "Input: \"Build a Vue 3 interview question tracker...\"  →  Output: 7 spec artifacts  •  34 tasks  •  61 tests  •  Deployed app",
            15, WHITE, False, PP_ALIGN.CENTER)


# ─── Slide 3: Problem & Solution ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Problem & Solution", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

# Problem box
add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0xFE, 0xF2, 0xF2), 0.03)
add_textbox(slide, Inches(1.2), Inches(2.0), Inches(4.8), Inches(0.6),
            "The Problem", 24, RGBColor(0xDC, 0x26, 0x26), True)
add_bullet_list(slide, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3.2), [
    "• Interview questions are scattered across notes, bookmarks, and memory",
    "• No easy way to organize questions by technology/category",
    "• Questions get lost when switching devices or clearing browser data",
    "• No quick way to filter and review domain-specific questions",
], 16, DARK)

# Solution box
add_shape(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(0xF0, 0xFD, 0xF4), 0.03)
add_textbox(slide, Inches(7.4), Inches(2.0), Inches(4.8), Inches(0.6),
            "Our Solution", 24, GREEN, True)
add_bullet_list(slide, Inches(7.4), Inches(2.7), Inches(4.8), Inches(3.2), [
    "• Centralized web app to track all interview questions",
    "• 6 built-in categories: Frontend, Backend, Vue, React, .NET, Others",
    "• Full CRUD: Add, View, Edit, Delete questions with notes",
    "• Persistent localStorage — data survives refreshes",
    "• Deployed to GitHub Pages — accessible anywhere",
], 16, DARK)


# ─── Slide 3: Tech Stack ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Technology Stack", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

techs = [
    ("Vue 3", "Composition API, reactive state, component architecture", RGBColor(0x05, 0x96, 0x69)),
    ("Vite 5", "Fast dev server, HMR, optimized production builds", ORANGE),
    ("Tailwind CSS 4", "Utility-first styling, responsive design, color-coded badges", CYAN),
    ("localStorage", "Client-side JSON persistence, no backend needed", PURPLE),
    ("Vitest", "61 unit & component tests, happy-dom environment", GREEN),
    ("GitHub Actions", "Automated CI/CD pipeline, deploy to GitHub Pages", DARK),
]

for i, (name, desc, color) in enumerate(techs):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 1.8)
    add_shape(slide, x, y, Inches(5.5), Inches(1.5), LIGHT_BG, 0.03)
    add_textbox(slide, x + Inches(0.4), y + Inches(0.2), Inches(4.8), Inches(0.5), name, 22, color, True)
    add_textbox(slide, x + Inches(0.4), y + Inches(0.8), Inches(4.8), Inches(0.6), desc, 15, MEDIUM)


# ─── Slide 4: Architecture ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Application Architecture", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

# Architecture diagram as boxes
# App.vue
add_shape(slide, Inches(4.5), Inches(1.8), Inches(4.3), Inches(0.9), PRIMARY, 0.05)
add_textbox(slide, Inches(4.5), Inches(1.9), Inches(4.3), Inches(0.8), "App.vue (Root)", 20, WHITE, True, PP_ALIGN.CENTER)

# Components layer
components = ["QuestionForm", "QuestionList", "QuestionCard", "CategoryFilter", "EmptyState"]
for i, comp in enumerate(components):
    x = Inches(0.5 + i * 2.5)
    add_shape(slide, x, Inches(3.3), Inches(2.2), Inches(0.8), RGBColor(0xDB, 0xEA, 0xFE), 0.05)
    add_textbox(slide, x, Inches(3.35), Inches(2.2), Inches(0.7), comp, 13, PRIMARY, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(3.0), Inches(2), Inches(0.3), "Components", 12, MEDIUM, True)

# Composables layer
composables = ["useQuestions.js", "useLocalStorage.js"]
for i, comp in enumerate(composables):
    x = Inches(3.2 + i * 3.5)
    add_shape(slide, x, Inches(4.8), Inches(3), Inches(0.8), RGBColor(0xD1, 0xFA, 0xE5), 0.05)
    add_textbox(slide, x, Inches(4.85), Inches(3), Inches(0.7), comp, 15, GREEN, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(4.9), Inches(2.5), Inches(0.3), "Composables (Logic)", 12, MEDIUM, True)

# Storage layer
add_shape(slide, Inches(3.8), Inches(6.2), Inches(5.5), Inches(0.8), RGBColor(0xF3, 0xE8, 0xFF), 0.05)
add_textbox(slide, Inches(3.8), Inches(6.25), Inches(5.5), Inches(0.7), "localStorage (Browser Persistence)", 16, PURPLE, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.5), Inches(6.3), Inches(2.5), Inches(0.3), "Storage Layer", 12, MEDIUM, True)


# ─── Slide 5: Data Model ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Data Model", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

# Question entity card
add_shape(slide, Inches(0.8), Inches(1.8), Inches(6), Inches(5), LIGHT_BG, 0.02)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(5), Inches(0.6), "Question Entity", 24, PRIMARY, True)

fields = [
    "id          string    UUID v4 (auto-generated, immutable)",
    "text        string    Question text (required, trimmed)",
    "category    string    One of 6 predefined values (required)",
    "notes       string    Optional answer/notes field",
    "createdAt   string    ISO 8601 timestamp (auto-set)",
]
add_textbox(slide, Inches(1.2), Inches(2.6), Inches(5.2), Inches(0.4), "Field           Type        Description", 13, MEDIUM, True)
for i, field in enumerate(fields):
    add_textbox(slide, Inches(1.2), Inches(3.1 + i * 0.5), Inches(5.2), Inches(0.4), field, 13, DARK)

# Categories card
add_shape(slide, Inches(7.5), Inches(1.8), Inches(5), Inches(5), LIGHT_BG, 0.02)
add_textbox(slide, Inches(7.9), Inches(1.9), Inches(4.2), Inches(0.6), "Categories (Fixed Enum)", 24, PRIMARY, True)

categories = [
    ("Frontend", RGBColor(0x1E, 0x40, 0xAF)),
    ("Backend", GREEN),
    ("Vue", RGBColor(0x05, 0x96, 0x69)),
    ("React", CYAN),
    (".NET", PURPLE),
    ("Others", MEDIUM),
]
for i, (cat, color) in enumerate(categories):
    y = Inches(2.7 + i * 0.65)
    add_shape(slide, Inches(8.2), y, Inches(3.5), Inches(0.5), WHITE, 0.1)
    add_textbox(slide, Inches(8.5), y + Inches(0.05), Inches(3), Inches(0.4), cat, 17, color, True, PP_ALIGN.LEFT)


# ─── Slide 6: User Stories ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "User Stories & Priorities", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

stories = [
    ("P1", "US1: Add a New Question", "Users can add questions with text, category, and optional notes", RGBColor(0xDC, 0x26, 0x26)),
    ("P1", "US2: View All Questions", "Card-based list sorted newest-first with category badges", RGBColor(0xDC, 0x26, 0x26)),
    ("P2", "US3: Edit an Existing Question", "Modify text, category, or notes via modal form", ORANGE),
    ("P2", "US4: Delete a Question", "Remove questions with confirmation dialog", ORANGE),
    ("P2", "US6: Data Persistence", "All data persists in localStorage across sessions", ORANGE),
    ("P3", "US5: Filter by Category", "Filter question list by category, clear to show all", CYAN),
    ("P3", "US7: Deploy to GitHub Pages", "Automated CI/CD via GitHub Actions", CYAN),
]

for i, (priority, title, desc, color) in enumerate(stories):
    y = Inches(1.7 + i * 0.78)
    # Priority badge
    add_shape(slide, Inches(0.8), y, Inches(0.8), Inches(0.55), color, 0.1)
    add_textbox(slide, Inches(0.8), y + Inches(0.05), Inches(0.8), Inches(0.45), priority, 14, WHITE, True, PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, Inches(1.8), y + Inches(0.02), Inches(4), Inches(0.5), title, 17, DARK, True)
    # Description
    add_textbox(slide, Inches(6), y + Inches(0.02), Inches(6.5), Inches(0.5), desc, 15, MEDIUM)


# ─── Slide 7: Key Features Demo ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Key Features", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

features = [
    ("CRUD Operations", "Full Create, Read, Update, Delete lifecycle\nwith form validation and confirmation dialogs", "📝"),
    ("Category System", "6 color-coded categories with badge indicators\nFrontend, Backend, Vue, React, .NET, Others", "🏷️"),
    ("Smart Filtering", "Filter questions by any category\nEmpty state messages for filtered views", "🔍"),
    ("Persistent Storage", "localStorage with JSON serialization\nData survives refresh, close, restart", "💾"),
    ("Responsive Design", "Mobile-first layout (320px minimum)\nTailwind CSS utility classes", "📱"),
    ("Automated Deploy", "GitHub Actions CI/CD pipeline\nBuild & deploy to GitHub Pages", "🚀"),
]

for i, (title, desc, icon) in enumerate(features):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.8 + row * 1.8)
    add_shape(slide, x, y, Inches(5.5), Inches(1.5), LIGHT_BG, 0.03)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.15), Inches(0.6), Inches(0.5), icon, 28, DARK, False, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.0), y + Inches(0.1), Inches(4.2), Inches(0.4), title, 19, PRIMARY, True)
    txBox = slide.shapes.add_textbox(x + Inches(1.0), y + Inches(0.6), Inches(4.2), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line_idx, line in enumerate(desc.split("\n")):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.color.rgb = MEDIUM


# ─── Slide 8: Implementation Phases ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Implementation Phases (34 Tasks)", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

phases = [
    ("Phase 1", "Setup", "Scaffold Vue 3 + Vite + Tailwind CSS project", "4 tasks", GREEN),
    ("Phase 2", "Foundation", "Categories, localStorage composable, useQuestions", "3 tasks", GREEN),
    ("Phase 3", "US1+US2 (P1)", "Add & View questions — MVP delivery", "5 tasks", GREEN),
    ("Phase 4", "US3 (P2)", "Edit questions via modal form", "3 tasks", GREEN),
    ("Phase 5", "US4 (P2)", "Delete with confirmation dialog", "2 tasks", GREEN),
    ("Phase 6", "US6 (P2)", "Verify localStorage persistence", "1 task", GREEN),
    ("Phase 7", "US5 (P3)", "Category filtering + empty states", "3 tasks", GREEN),
    ("Phase 8", "US7 (P3)", "GitHub Actions deployment", "2 tasks", GREEN),
    ("Phase 9", "Polish", "Responsive design, edge cases, validation", "3 tasks", GREEN),
    ("Phase 10", "Testing", "61 tests across 7 test files", "8 tasks", GREEN),
]

for i, (phase, name, desc, tasks, color) in enumerate(phases):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.7 + row * 1.1)

    # Phase badge
    add_shape(slide, x, y, Inches(5.5), Inches(0.9), LIGHT_BG, 0.03)
    add_shape(slide, x, y, Inches(1.2), Inches(0.9), PRIMARY, 0.03)
    add_textbox(slide, x, y + Inches(0.05), Inches(1.2), Inches(0.35), phase, 11, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(0.4), Inches(1.2), Inches(0.35), tasks, 10, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(1.4), y + Inches(0.08), Inches(3.8), Inches(0.35), name, 16, DARK, True)
    add_textbox(slide, x + Inches(1.4), y + Inches(0.45), Inches(3.8), Inches(0.4), desc, 13, MEDIUM)


# ─── Slide 9: Testing Summary ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "Testing & Quality", 36, PRIMARY, True)
add_shape(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.05), ACCENT)

# Big number
add_shape(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(2.5), PRIMARY, 0.03)
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(3.5), Inches(1.2), "61", 72, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(3.2), Inches(3.5), Inches(0.6), "Total Tests", 22, RGBColor(0xBF, 0xDB, 0xFE), False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(3.7), Inches(3.5), Inches(0.5), "All Passing ✓", 18, RGBColor(0x86, 0xEF, 0xAC), False, PP_ALIGN.CENTER)

# Test file breakdown
test_files = [
    ("useQuestions.test.js", "19 tests", "CRUD, validation, filtering, sorting, HTML stripping"),
    ("QuestionForm.test.js", "12 tests", "Add/edit modes, validation, emit save/cancel"),
    ("QuestionCard.test.js", "7 tests", "Display, notes visibility, edit/delete emit"),
    ("QuestionList.test.js", "7 tests", "Empty/populated states, event bubbling"),
    ("CategoryFilter.test.js", "6 tests", "Rendering, active state, filter changes"),
    ("categories.test.js", "5 tests", "Constants and badge class helper"),
    ("useLocalStorage.test.js", "5 tests", "Load, save, corrupted JSON, reactivity"),
]

for i, (fname, count, desc) in enumerate(test_files):
    y = Inches(1.8 + i * 0.72)
    add_shape(slide, Inches(5), y, Inches(7.5), Inches(0.6), LIGHT_BG, 0.05)
    add_textbox(slide, Inches(5.2), y + Inches(0.05), Inches(2.5), Inches(0.5), fname, 13, DARK, True)
    add_shape(slide, Inches(7.8), y + Inches(0.08), Inches(1), Inches(0.4), GREEN, 0.15)
    add_textbox(slide, Inches(7.8), y + Inches(0.08), Inches(1), Inches(0.4), count, 12, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(9), y + Inches(0.08), Inches(3.5), Inches(0.5), desc, 12, MEDIUM)


# ─── Slide 10: Summary & Next Steps ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, PRIMARY)

add_textbox(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
            "Summary & Next Steps", 40, WHITE, True, PP_ALIGN.CENTER)

# Delivered column
add_shape(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), RGBColor(0x16, 0x35, 0x8A), 0.03)
add_textbox(slide, Inches(1.2), Inches(2.4), Inches(4.8), Inches(0.6), "✅  What We Delivered", 24, WHITE, True)
add_bullet_list(slide, Inches(1.2), Inches(3.1), Inches(4.8), Inches(3.2), [
    "• Full CRUD operations with validation",
    "• 6 color-coded categories with filtering",
    "• Persistent localStorage across sessions",
    "• Responsive design (mobile-friendly)",
    "• 61 passing tests (unit + component)",
    "• Automated GitHub Pages deployment",
    "• Single-page Vue 3 app, no backend needed",
], 16, RGBColor(0xBF, 0xDB, 0xFE))

# Next steps column
add_shape(slide, Inches(7), Inches(2.2), Inches(5.5), Inches(4.5), RGBColor(0x16, 0x35, 0x8A), 0.03)
add_textbox(slide, Inches(7.4), Inches(2.4), Inches(4.8), Inches(0.6), "🔮  Potential Enhancements", 24, WHITE, True)
add_bullet_list(slide, Inches(7.4), Inches(3.1), Inches(4.8), Inches(3.2), [
    "• Search across question text and notes",
    "• Export/import questions (JSON/CSV)",
    "• Custom user-defined categories",
    "• Drag-and-drop question reordering",
    "• Dark mode toggle",
    "• Cloud sync with user accounts",
    "• Spaced repetition for review",
], 16, RGBColor(0xBF, 0xDB, 0xFE))

add_textbox(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.5),
            "Thank you!  •  Questions?", 20, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER)


# Save
prs.save("Interview_Question_Tracker_Summary_v2.pptx")
print("Presentation saved: Interview_Question_Tracker_Summary.pptx")
